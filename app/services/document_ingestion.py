"""
Document ingestion service.

Uses individual format-specific libraries instead of markitdown to avoid the
onnxruntime DLL initialisation failure on Python 3.14 + Windows (triggered by
markitdown's required dependency `magika` which imports onnxruntime at module
load time).

All libraries used here were already installed as markitdown transitive deps:
  pdfplumber  — PDF text extraction
  mammoth     — DOCX → Markdown
  python-pptx — PPTX text extraction
  beautifulsoup4 / lxml — HTML text extraction
"""

import csv
import io
import os


# ── Format-specific converters ────────────────────────────────────────────────

def _pdf_to_markdown(file_bytes: bytes) -> str:
    import pdfplumber
    parts = []
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text and text.strip():
                parts.append(text.strip())
    return "\n\n".join(parts) or "(no text could be extracted from this PDF)"


def _docx_to_markdown(file_bytes: bytes) -> str:
    import mammoth
    result = mammoth.convert_to_markdown(io.BytesIO(file_bytes))
    return result.value.strip() or "(no text could be extracted from this DOCX)"


def _pptx_to_markdown(file_bytes: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(file_bytes))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        texts = [
            shape.text.strip()
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        ]
        if texts:
            parts.append(f"## Slide {i}\n\n" + "\n\n".join(texts))
    return "\n\n".join(parts) or "(no text could be extracted from this PPTX)"


def _html_to_markdown(file_bytes: bytes) -> str:
    from bs4 import BeautifulSoup
    soup = BeautifulSoup(file_bytes, "lxml")
    for tag in soup(["script", "style"]):
        tag.decompose()
    return soup.get_text(separator="\n").strip() or "(no text could be extracted from this HTML)"


def _csv_to_markdown(file_bytes: bytes) -> str:
    text = file_bytes.decode("utf-8", errors="replace")
    rows = list(csv.reader(text.splitlines()))
    if not rows:
        return "(empty CSV)"
    header = rows[0]
    sep = "| " + " | ".join(["---"] * len(header)) + " |"
    lines = ["| " + " | ".join(header) + " |", sep]
    for row in rows[1:]:
        # Pad short rows to match header length
        padded = row + [""] * (len(header) - len(row))
        lines.append("| " + " | ".join(padded[: len(header)]) + " |")
    return "\n".join(lines)


_CONVERTERS = {
    ".pdf":  _pdf_to_markdown,
    ".docx": _docx_to_markdown,
    ".pptx": _pptx_to_markdown,
    ".html": _html_to_markdown,
    ".htm":  _html_to_markdown,
    ".csv":  _csv_to_markdown,
    ".tsv":  _csv_to_markdown,
}


# ── Public API ────────────────────────────────────────────────────────────────

def ingest_document(file_bytes: bytes, filename: str) -> dict:
    """
    Convert file bytes to Markdown.

    Supported formats: PDF, DOCX, PPTX, HTML, CSV/TSV.
    All other files are decoded as UTF-8 text (works for .txt, .md, .json,
    .xml, and any other plain-text format).
    """
    ext = os.path.splitext(filename)[1].lower()
    converter = _CONVERTERS.get(ext)

    if converter:
        markdown = converter(file_bytes)
    else:
        # Plain-text fallback (txt, md, json, xml, …)
        try:
            markdown = file_bytes.decode("utf-8", errors="replace")
        except Exception:
            markdown = "(binary file — no text could be extracted)"

    return {
        "markdown": markdown,
        "original_bytes": len(file_bytes),
        "filename": filename,
    }
